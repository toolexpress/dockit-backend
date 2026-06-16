"""
DocKit Backend — FastAPI
Handles: PDF↔Word, PDF↔PPT, PDF↔Excel conversions via LibreOffice + pdf2docx
Deploy on Railway / Render / any cloud VPS

SECURITY: Dual-layer CORS (Website origin + Mobile app null/capacitor origin)
         + X-App-Token verification for request authenticity
"""

import os
import uuid
import shutil
import subprocess
import tempfile
import secrets
import time
from pathlib import Path

from fastapi import FastAPI, File, UploadFile, HTTPException, Request, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pdf2docx import Converter as PDF2DocxConverter

app = FastAPI(title="DocKit API", version="1.0.0")

# ── SECURITY CONFIG ───────────────────────────────────────────────────────
# App token: shared secret embedded in website JS + mobile app config.
# This is NOT a substitute for real auth (it's visible in client code),
# but it blocks casual scrapers/bots that don't bother inspecting JS/app,
# and lets you tell legit app/website traffic apart from random hits.
APP_TOKEN = os.environ.get("APP_TOKEN", "dockit-secret-2026-change-me")

# Origins allowed to call the API from a browser (Website layer)
WEB_ORIGINS = [
    "https://pdfdockit.toolexpress.online",
    "https://www.pdfdockit.toolexpress.online",
]

# Mobile apps (Capacitor) send Origin as "null", "capacitor://localhost",
# "http://localhost", or no Origin header at all — never a real domain.
MOBILE_ORIGINS = [
    "capacitor://localhost",
    "http://localhost",
    "https://localhost",
    "null",
]

ALL_ORIGINS = WEB_ORIGINS + MOBILE_ORIGINS

# ── CORS — Dual layer: website origins + mobile app origins ─────────────────
app.add_middleware(
    CORSMiddleware,
    allow_origins=ALL_ORIGINS,
    allow_credentials=False,  # no cookies used, keep False with explicit origins
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)


# ── APP TOKEN VERIFICATION MIDDLEWARE ────────────────────────────────────────
# Every request must carry X-App-Token header matching APP_TOKEN.
# Health check + docs are exempt so monitoring/uptime services still work.
EXEMPT_PATHS = {"/api/health", "/docs", "/openapi.json", "/redoc"}

@app.middleware("http")
async def verify_app_token(request: Request, call_next):
    if request.method == "OPTIONS" or request.url.path in EXEMPT_PATHS:
        return await call_next(request)

    token = request.headers.get("x-app-token")
    if not token or not secrets.compare_digest(token, APP_TOKEN):
        return JSONResponse(
            status_code=401,
            content={"detail": "Unauthorized: missing or invalid X-App-Token"},
        )
    return await call_next(request)


# ── SIMPLE RATE LIMIT (per-IP, in-memory) ────────────────────────────────────
# Lightweight abuse guard. For production scale, swap for Redis-backed limiter.
_rate_buckets: dict[str, list[float]] = {}
RATE_LIMIT = 20          # max requests
RATE_WINDOW = 60         # per 60 seconds

@app.middleware("http")
async def rate_limit(request: Request, call_next):
    if request.url.path in EXEMPT_PATHS or request.method == "OPTIONS":
        return await call_next(request)

    client_ip = request.client.host if request.client else "unknown"
    now = time.time()
    bucket = _rate_buckets.setdefault(client_ip, [])
    # drop timestamps outside the window
    bucket[:] = [t for t in bucket if now - t < RATE_WINDOW]
    if len(bucket) >= RATE_LIMIT:
        return JSONResponse(
            status_code=429,
            content={"detail": "Too many requests. Please slow down."},
        )
    bucket.append(now)
    return await call_next(request)


TEMP_DIR = Path(tempfile.gettempdir()) / "dockit"
TEMP_DIR.mkdir(exist_ok=True)


# ── HELPERS ──────────────────────────────────────────────────────────────────

def save_upload(file: UploadFile, suffix: str) -> Path:
    """Save uploaded file to temp dir, return path."""
    dest = TEMP_DIR / f"{uuid.uuid4().hex}{suffix}"
    with open(dest, "wb") as f:
        shutil.copyfileobj(file.file, f)
    return dest


def libreoffice_convert(input_path: Path, output_format: str, output_dir: Path) -> Path:
    """Run LibreOffice headless conversion, return output file path."""
    cmd = [
        "libreoffice", "--headless", "--convert-to", output_format,
        "--outdir", str(output_dir), str(input_path)
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice error: {result.stderr}")

    # LibreOffice outputs file with same stem, different extension
    out_file = output_dir / (input_path.stem + "." + output_format.split(":")[0])
    if not out_file.exists():
        # Try finding any newly created file in output_dir
        candidates = list(output_dir.glob(f"*.{output_format.split(':')[0]}"))
        if not candidates:
            raise RuntimeError("Output file not found after conversion")
        out_file = candidates[0]
    return out_file


def cleanup(*paths):
    for p in paths:
        try:
            if p and Path(p).exists():
                os.remove(p)
        except Exception:
            pass


# ── HEALTH CHECK ─────────────────────────────────────────────────────────────

@app.get("/api/health")
def health():
    """Frontend pings this to check if server is online. No token required."""
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    return {
        "status": "online",
        "libreoffice": bool(lo),
        "libreoffice_path": lo
    }


# ── 1. PDF → WORD (.docx) ────────────────────────────────────────────────────

@app.post("/api/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "PDF file required")

    src = save_upload(file, ".pdf")
    out = TEMP_DIR / f"{uuid.uuid4().hex}.docx"
    try:
        cv = PDF2DocxConverter(str(src))
        cv.convert(str(out), start=0, end=None)
        cv.close()
        return FileResponse(
            str(out), media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=Path(file.filename).stem + ".docx",
            background=None
        )
    except Exception as e:
        cleanup(src, out)
        raise HTTPException(500, f"Conversion failed: {e}")
    finally:
        cleanup(src)


# ── 2. PDF → POWERPOINT (.pptx) ──────────────────────────────────────────────

@app.post("/api/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    """
    Convert PDF pages → PPTX (each page = one slide image).
    Uses pdf2image + python-pptx.
    """
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "PDF file required")

    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import io
    except ImportError as e:
        raise HTTPException(500, f"Missing dependency: {e}. Install pdf2image and python-pptx.")

    src = save_upload(file, ".pdf")
    out = TEMP_DIR / f"{uuid.uuid4().hex}.pptx"
    try:
        images = convert_from_path(str(src), dpi=150)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]  # blank

        for img in images:
            slide = prs.slides.add_slide(blank_layout)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=85)
            buf.seek(0)
            slide.shapes.add_picture(buf, 0, 0, prs.slide_width, prs.slide_height)

        prs.save(str(out))
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=Path(file.filename).stem + ".pptx"
        )
    except Exception as e:
        cleanup(src, out)
        raise HTTPException(500, f"Conversion failed: {e}")
    finally:
        cleanup(src)


# ── 3. PDF → EXCEL (.xlsx) ───────────────────────────────────────────────────

@app.post("/api/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Extract text/tables from PDF into Excel using pdfplumber."""
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "PDF file required")

    try:
        import pdfplumber
        import openpyxl
    except ImportError as e:
        raise HTTPException(500, f"Missing dependency: {e}")

    src = save_upload(file, ".pdf")
    out = TEMP_DIR / f"{uuid.uuid4().hex}.xlsx"
    try:
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default sheet

        with pdfplumber.open(str(src)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                ws = wb.create_sheet(title=f"Page {i}")
                tables = page.extract_tables()
                if tables:
                    row_idx = 1
                    for table in tables:
                        for row in table:
                            for col_idx, cell in enumerate(row, 1):
                                ws.cell(row=row_idx, column=col_idx, value=cell or "")
                            row_idx += 1
                        row_idx += 1  # blank row between tables
                else:
                    # No tables — dump raw text
                    text = page.extract_text() or ""
                    for row_idx, line in enumerate(text.split("\n"), 1):
                        ws.cell(row=row_idx, column=1, value=line)

        wb.save(str(out))
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=Path(file.filename).stem + ".xlsx"
        )
    except Exception as e:
        cleanup(src, out)
        raise HTTPException(500, f"Conversion failed: {e}")
    finally:
        cleanup(src)


# ── 4. WORD → PDF ─────────────────────────────────────────────────────────────

@app.post("/api/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    fname = file.filename.lower()
    if not (fname.endswith(".doc") or fname.endswith(".docx")):
        raise HTTPException(400, ".doc or .docx file required")

    suffix = ".docx" if fname.endswith(".docx") else ".doc"
    src = save_upload(file, suffix)
    out_dir = TEMP_DIR
    try:
        out = libreoffice_convert(src, "pdf", out_dir)
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=Path(file.filename).stem + ".pdf"
        )
    except Exception as e:
        raise HTTPException(500, f"Conversion failed: {e}")
    finally:
        cleanup(src)


# ── 5. POWERPOINT → PDF ──────────────────────────────────────────────────────

@app.post("/api/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    fname = file.filename.lower()
    if not (fname.endswith(".ppt") or fname.endswith(".pptx")):
        raise HTTPException(400, ".ppt or .pptx file required")

    suffix = ".pptx" if fname.endswith(".pptx") else ".ppt"
    src = save_upload(file, suffix)
    try:
        out = libreoffice_convert(src, "pdf", TEMP_DIR)
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=Path(file.filename).stem + ".pdf"
        )
    except Exception as e:
        raise HTTPException(500, f"Conversion failed: {e}")
    finally:
        cleanup(src)


# ── 6. EXCEL → PDF ───────────────────────────────────────────────────────────

@app.post("/api/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    fname = file.filename.lower()
    if not (fname.endswith(".xls") or fname.endswith(".xlsx")):
        raise HTTPException(400, ".xls or .xlsx file required")

    suffix = ".xlsx" if fname.endswith(".xlsx") else ".xls"
    src = save_upload(file, suffix)
    try:
        out = libreoffice_convert(src, "pdf", TEMP_DIR)
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=Path(file.filename).stem + ".pdf"
        )
    except Exception as e:
        raise HTTPException(500, f"Conversion failed: {e}")
    finally:
        cleanup(src)


# ── OCR PDF ──────────────────────────────────────────────────────────────────

@app.post("/api/ocr-pdf")
async def ocr_pdf(file: UploadFile = File(...)):
    """
    Run OCR on a scanned PDF using ocrmypdf (Tesseract under the hood).
    Returns a searchable PDF with invisible text layer.
    """
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "PDF file required")

    src = save_upload(file, ".pdf")
    out = TEMP_DIR / f"{uuid.uuid4().hex}_ocr.pdf"
    try:
        cmd = [
            "ocrmypdf",
            "--skip-text",          # skip pages that already have text
            "--optimize", "1",
            "--output-type", "pdf",
            str(src), str(out)
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        if result.returncode not in (0, 6):  # 6 = already has text (ok)
            raise RuntimeError(result.stderr or "OCR failed")
        # If output wasn't created (all pages skipped), return original
        if not out.exists():
            out = src
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=Path(file.filename).stem + "_ocr.pdf"
        )
    except Exception as e:
        cleanup(src, out)
        raise HTTPException(500, f"OCR failed: {e}")
    finally:
        cleanup(src)


# ── RUN ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("main:app", host="0.0.0.0", port=port, reload=False)
